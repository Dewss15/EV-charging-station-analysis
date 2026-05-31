"""
CE634 Data Mining — Premium Viva Presentation Generator
Design: McKinsey/BCG × Apple aesthetic
Palette: Deep Navy | Electric Blue | Teal | White | Light Grey
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── paths ──────────────────────────────────────────────────────────────────────
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMG = {
    "correlation": os.path.join(BASE_DIR, "correlation.png"),
    "euclidean":   os.path.join(BASE_DIR, "euclidean.png"),
    "confusion":   os.path.join(BASE_DIR, "confusion_matrix.png"),
    "feature_imp": os.path.join(BASE_DIR, "feature_imp.png"),
    "class_dist":  os.path.join(BASE_DIR, "class.png"),
    "dashboard":   os.path.join(BASE_DIR, "final_dashboard.png"),
}
OUTPUT = os.path.join(BASE_DIR, "CE634_EV_Charging_Presentation.pptx")

# ── colour palette ─────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # Deep navy background
NAVY_LIGHT  = RGBColor(0x12, 0x23, 0x52)   # Lighter navy for cards
BLUE        = RGBColor(0x00, 0x8B, 0xFF)   # Electric blue accent
TEAL        = RGBColor(0x00, 0xC9, 0xB1)   # Teal accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # Pure white
OFFWHITE    = RGBColor(0xF0, 0xF4, 0xFF)   # Off-white text
GREY        = RGBColor(0xA8, 0xB8, 0xD8)   # Light grey supporting text
CARD_BG     = RGBColor(0x16, 0x2A, 0x5A)   # Card background
DIVIDER     = RGBColor(0x00, 0x6B, 0xC4)   # Divider line
HIGHLIGHT   = RGBColor(0xFF, 0xC1, 0x07)   # Amber highlight
GREEN_STAT  = RGBColor(0x00, 0xE5, 0x96)   # Green for good stats
RED_WARN    = RGBColor(0xFF, 0x4D, 0x6D)   # Red for warnings
ORANGE      = RGBColor(0xFF, 0x8C, 0x42)   # Orange accent

# ── slide dimensions (16:9 widescreen) ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════

def new_prs():
    """Create a new blank widescreen presentation."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    """Add a completely blank slide."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_slide_bg(slide, color: RGBColor):
    """Fill the entire slide background with a solid colour."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color: RGBColor = None,
         line_color: RGBColor = None, line_width_pt=0, transparency=0):
    """Add a rectangle shape."""
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width_pt > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def rounded_rect(slide, left, top, width, height, fill_color: RGBColor,
                 line_color=None, line_width_pt=0, corner_radius=0.08):
    """Add a rounded-corner rectangle."""
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    # set corner radius
    try:
        adj = shape.adjustments
        adj[0] = corner_radius
    except Exception:
        pass
    return shape


def txbox(slide, text, left, top, width, height,
          font_size=18, bold=False, color: RGBColor = WHITE,
          align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    """Add a text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_image(slide, img_path, left, top, width, height=None):
    """Add an image; if height is None, preserve aspect ratio."""
    if not os.path.exists(img_path):
        print(f"  [WARN] Image not found: {img_path}")
        return None
    if height:
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        pic = slide.shapes.add_picture(img_path, left, top, width)
    return pic


def divider_line(slide, left, top, width, color: RGBColor = BLUE, thickness_pt=1.5):
    """Draw a thin horizontal line."""
    from pptx.util import Pt as _Pt
    connector = slide.shapes.add_shape(1, left, top, width, Pt(thickness_pt))
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def slide_number_tag(slide, number, total, color=GREY):
    """Add a small slide number in the bottom-right corner."""
    txbox(slide, f"{number} / {total}",
          SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.4),
          Inches(0.8), Inches(0.3),
          font_size=9, color=color, align=PP_ALIGN.RIGHT)


def accent_bar(slide, color=TEAL, width=Inches(0.05)):
    """Left-side accent bar spanning full height."""
    rect(slide, 0, 0, width, SLIDE_H, fill_color=color)


def section_pill(slide, label, left, top, bg=BLUE, text_color=WHITE, font_size=9):
    """Small pill-shaped label."""
    w = Inches(1.6)
    h = Inches(0.28)
    r = rounded_rect(slide, left, top, w, h, fill_color=bg, corner_radius=0.4)
    txbox(slide, label.upper(), left, top + Pt(2), w, h,
          font_size=font_size, bold=True, color=text_color,
          align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
def build_title_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)

    # left dark panel (2/3 width)
    rect(slide, 0, 0, Inches(8.8), SLIDE_H, fill_color=NAVY)

    # right accent panel
    rect(slide, Inches(8.8), 0, Inches(4.53), SLIDE_H,
         fill_color=RGBColor(0x08, 0x14, 0x30))

    # electric-blue top bar
    rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill_color=BLUE)

    # teal accent bar on left edge
    rect(slide, 0, 0, Inches(0.06), SLIDE_H, fill_color=TEAL)

    # college & course tag
    txbox(slide, "PADRE CONCEICAO COLLEGE OF ENGINEERING",
          Inches(0.4), Inches(0.35), Inches(8.0), Inches(0.35),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    txbox(slide, "CE634 · Data Mining and Data Warehousing",
          Inches(0.4), Inches(0.72), Inches(8.0), Inches(0.35),
          font_size=10, color=GREY, font_name="Calibri")

    # thin divider under header
    divider_line(slide, Inches(0.4), Inches(1.12), Inches(8.2), color=BLUE)

    # main title
    txbox(slide, "Predictive Analysis of",
          Inches(0.4), Inches(1.3), Inches(8.2), Inches(0.65),
          font_size=28, bold=False, color=OFFWHITE, font_name="Calibri")

    txbox(slide, "EV Charging Demand",
          Inches(0.4), Inches(1.9), Inches(8.2), Inches(0.9),
          font_size=42, bold=True, color=WHITE, font_name="Calibri")

    txbox(slide, "& Grid Load Dynamics",
          Inches(0.4), Inches(2.72), Inches(8.2), Inches(0.75),
          font_size=36, bold=True, color=BLUE, font_name="Calibri")

    # subtitle
    txbox(slide, "Peak Load Risk Classification · Operational Clustering · Anomaly Detection",
          Inches(0.4), Inches(3.6), Inches(8.2), Inches(0.45),
          font_size=12, color=GREY, italic=True, font_name="Calibri")

    divider_line(slide, Inches(0.4), Inches(4.15), Inches(8.2), color=DIVIDER)

    # team members
    txbox(slide, "Dewpearl Gonsalves  ·  23CE114",
          Inches(0.4), Inches(4.3), Inches(5.0), Inches(0.32),
          font_size=11, color=OFFWHITE, font_name="Calibri")
    txbox(slide, "Shail Joshi  ·  23CE163",
          Inches(0.4), Inches(4.62), Inches(5.0), Inches(0.32),
          font_size=11, color=OFFWHITE, font_name="Calibri")
    txbox(slide, "Flyson Dias  ·  23CE118",
          Inches(0.4), Inches(4.94), Inches(5.0), Inches(0.32),
          font_size=11, color=OFFWHITE, font_name="Calibri")

    txbox(slide, "Academic Year 2025–2026  |  May 2026",
          Inches(0.4), Inches(5.42), Inches(5.0), Inches(0.32),
          font_size=9.5, color=GREY, font_name="Calibri")

    # right panel — stats banner
    txbox(slide, "AT A GLANCE",
          Inches(9.0), Inches(0.8), Inches(4.0), Inches(0.35),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    stats = [
        ("2,800", "Charging Records"),
        ("10",    "Dataset Attributes"),
        ("0",     "Missing Values"),
        ("45.18%","RF Accuracy"),
        ("84",    "Anomalies Flagged"),
        ("3",     "Operational Clusters"),
    ]
    card_top = Inches(1.22)
    for val, label in stats:
        rounded_rect(slide, Inches(9.1), card_top, Inches(3.9), Inches(0.77),
                     fill_color=CARD_BG, corner_radius=0.12)
        txbox(slide, val,
              Inches(9.2), card_top + Inches(0.06), Inches(1.5), Inches(0.38),
              font_size=20, bold=True, color=BLUE, font_name="Calibri")
        txbox(slide, label,
              Inches(10.55), card_top + Inches(0.22), Inches(2.4), Inches(0.28),
              font_size=9, color=GREY, font_name="Calibri")
        card_top += Inches(0.84)

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 2: AGENDA ───────────────────────────────────────────────────────────
def build_agenda_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    txbox(slide, "AGENDA", Inches(0.25), Inches(0.22), Inches(3.0), Inches(0.35),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")
    txbox(slide, "What We'll Cover Today",
          Inches(0.25), Inches(0.55), Inches(12.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.25), Inches(12.8), color=BLUE)

    items = [
        ("01", "Dataset & Problem",          "2,800 EV records · 10 features · peak load risk target"),
        ("02", "Preprocessing Pipeline",     "Cleaning · Engineering · Encoding · Scaling"),
        ("03", "Similarity Analysis",         "Pearson correlation · Euclidean distance across city zones"),
        ("04", "Classification (RF)",         "Random Forest · SMOTE · 45.18% accuracy · 3-class"),
        ("05", "Feature Importance",          "Top predictors of grid stress events"),
        ("06", "Clustering & Anomalies",      "K-Means (k=3) · Isolation Forest (84 flagged events)"),
        ("07", "Dashboard & Conclusions",     "Executive summary · Findings · Future work"),
    ]

    col1_items = items[:4]
    col2_items = items[4:]

    for i, (num, title, sub) in enumerate(col1_items):
        top = Inches(1.5) + i * Inches(1.35)
        rounded_rect(slide, Inches(0.25), top, Inches(6.2), Inches(1.18),
                     fill_color=CARD_BG, corner_radius=0.08)
        txbox(slide, num, Inches(0.38), top + Inches(0.12), Inches(0.65), Inches(0.55),
              font_size=22, bold=True, color=BLUE, font_name="Calibri")
        txbox(slide, title, Inches(1.05), top + Inches(0.1), Inches(5.1), Inches(0.42),
              font_size=13, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, sub, Inches(1.05), top + Inches(0.52), Inches(5.1), Inches(0.55),
              font_size=9, color=GREY, font_name="Calibri")

    for i, (num, title, sub) in enumerate(col2_items):
        top = Inches(1.5) + i * Inches(1.35)
        rounded_rect(slide, Inches(6.9), top, Inches(6.2), Inches(1.18),
                     fill_color=CARD_BG, corner_radius=0.08)
        txbox(slide, num, Inches(7.03), top + Inches(0.12), Inches(0.65), Inches(0.55),
              font_size=22, bold=True, color=TEAL, font_name="Calibri")
        txbox(slide, title, Inches(7.7), top + Inches(0.1), Inches(5.1), Inches(0.42),
              font_size=13, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, sub, Inches(7.7), top + Inches(0.52), Inches(5.1), Inches(0.55),
              font_size=9, color=GREY, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 3: DATASET ──────────────────────────────────────────────────────────
def build_dataset_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "01 · Dataset", Inches(0.25), Inches(0.2))
    txbox(slide, "Dataset at a Glance",
          Inches(0.25), Inches(0.56), Inches(8.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    # Big 4 stat cards
    stat_cards = [
        ("2,800", "Charging\nRecords", BLUE),
        ("10",    "Dataset\nAttributes", TEAL),
        ("0",     "Missing\nValues", GREEN_STAT),
        ("0",     "Duplicate\nRecords", GREEN_STAT),
    ]
    card_w = Inches(2.85)
    for i, (val, lbl, clr) in enumerate(stat_cards):
        cx = Inches(0.25) + i * (card_w + Inches(0.2))
        rounded_rect(slide, cx, Inches(1.45), card_w, Inches(1.6),
                     fill_color=CARD_BG, corner_radius=0.1)
        # accent top line on card
        rect(slide, cx, Inches(1.45), card_w, Inches(0.06), fill_color=clr)
        txbox(slide, val, cx + Inches(0.15), Inches(1.6), card_w - Inches(0.3), Inches(0.78),
              font_size=38, bold=True, color=clr, font_name="Calibri")
        txbox(slide, lbl, cx + Inches(0.15), Inches(2.35), card_w - Inches(0.3), Inches(0.6),
              font_size=10, color=GREY, font_name="Calibri")

    # Feature table header
    txbox(slide, "FEATURE OVERVIEW",
          Inches(0.25), Inches(3.25), Inches(5.0), Inches(0.32),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    features_left = [
        ("record_id",                    "Integer",     "Session identifier (excluded from model)"),
        ("date_time",                    "DateTime",    "Timestamp → hour_of_day, day_of_week"),
        ("city_zone",                    "Categorical", "North / South / East / West / Central"),
        ("station_type",                 "Categorical", "Fast / Slow / Ultra-fast / Normal"),
        ("vehicles_charged",             "Integer",     "No. of EVs per session"),
    ]
    features_right = [
        ("avg_charging_duration_min",    "Float",       "Mean session duration (minutes)"),
        ("energy_dispensed_kwh",         "Float",       "Total energy delivered (kWh)"),
        ("grid_load_mw",                 "Float",       "Grid draw during session (MW)"),
        ("renewable_energy_percent",     "Float",       "% renewable energy used"),
        ("peak_load_risk  ★ TARGET",     "Categorical", "Low / Medium / High"),
    ]

    def feat_row(feats, left_x, col_w):
        for i, (name, dtype, desc) in enumerate(feats):
            row_top = Inches(3.62) + i * Inches(0.66)
            bg = CARD_BG if i % 2 == 0 else RGBColor(0x0F, 0x1E, 0x48)
            rect(slide, left_x, row_top, col_w, Inches(0.6), fill_color=bg)
            txbox(slide, name, left_x + Inches(0.1), row_top + Inches(0.04),
                  Inches(1.9), Inches(0.3), font_size=9, bold=True, color=OFFWHITE, font_name="Calibri")
            txbox(slide, dtype, left_x + Inches(2.05), row_top + Inches(0.04),
                  Inches(0.9), Inches(0.3), font_size=8, color=TEAL, font_name="Calibri")
            txbox(slide, desc, left_x + Inches(0.1), row_top + Inches(0.32),
                  col_w - Inches(0.2), Inches(0.26), font_size=7.5, color=GREY, font_name="Calibri")

    # Col headers
    for cx, label in [(Inches(0.25), "Feature"), (Inches(6.75), "Feature")]:
        rect(slide, cx, Inches(3.55), Inches(6.35), Inches(0.3),
             fill_color=RGBColor(0x00, 0x6B, 0xC4))
        for fx, fl in [(cx + Inches(0.1), "FEATURE"), (cx + Inches(2.0), "TYPE"),
                       (cx + Inches(3.0), "DESCRIPTION")]:
            txbox(slide, fl, fx, Inches(3.58), Inches(1.8), Inches(0.26),
                  font_size=7.5, bold=True, color=WHITE, font_name="Calibri")

    feat_row(features_left,  Inches(0.25), Inches(6.35))
    feat_row(features_right, Inches(6.75), Inches(6.35))

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 4: PREPROCESSING ────────────────────────────────────────────────────
def build_preprocessing_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "02 · Preprocessing", Inches(0.25), Inches(0.2))
    txbox(slide, "8-Step Preprocessing Pipeline",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    txbox(slide, "Input: Raw CSV (2,800 × 10)  →  Output: Analysis-Ready Matrix (2,800 × 15)",
          Inches(0.25), Inches(1.35), Inches(12.0), Inches(0.38),
          font_size=10, color=GREY, italic=True, font_name="Calibri")

    steps = [
        ("1", "Null\nCheck",        "Confirmed 0\nmissing values",   TEAL),
        ("2", "Duplicate\nRemoval", "0 duplicates\nfound",           TEAL),
        ("3", "Datetime\nParsing",  "Extracted temporal\nfeatures",  BLUE),
        ("4", "Feature\nEngineering","hour_of_day\nday_of_week",     BLUE),
        ("5", "One-Hot\nEncoding",  "station_type,\ncity_zone",      ORANGE),
        ("6", "Target\nMapping",    "Low→0, Med→1\nHigh→2",          ORANGE),
        ("7", "StandardScaler",     "Z-score normalize\n5 features", HIGHLIGHT),
        ("8", "Drop\nrecord_id",    "Prevent data\nleakage",         RED_WARN),
    ]

    step_w = Inches(1.51)
    arrow_w = Inches(0.18)
    total_w = len(steps) * step_w + (len(steps) - 1) * arrow_w
    start_x = (SLIDE_W - total_w) / 2

    for i, (num, title, sub, clr) in enumerate(steps):
        cx = start_x + i * (step_w + arrow_w)
        cy = Inches(1.88)
        card_h = Inches(3.6)

        # card
        rounded_rect(slide, cx, cy, step_w, card_h, fill_color=CARD_BG, corner_radius=0.1)
        # top colour accent
        rect(slide, cx, cy, step_w, Inches(0.06), fill_color=clr)

        # step number circle (simulated with a rectangle)
        rounded_rect(slide, cx + Inches(0.48), cy + Inches(0.15),
                     Inches(0.52), Inches(0.52), fill_color=clr, corner_radius=0.5)
        txbox(slide, num,
              cx + Inches(0.48), cy + Inches(0.17), Inches(0.52), Inches(0.44),
              font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

        # title
        txbox(slide, title,
              cx + Inches(0.06), cy + Inches(0.78), step_w - Inches(0.12), Inches(0.75),
              font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

        # sub detail
        txbox(slide, sub,
              cx + Inches(0.06), cy + Inches(1.58), step_w - Inches(0.12), Inches(0.85),
              font_size=8, color=GREY, align=PP_ALIGN.CENTER, font_name="Calibri")

        # arrow between cards
        if i < len(steps) - 1:
            ax = cx + step_w + Inches(0.02)
            txbox(slide, "→",
                  ax, cy + Inches(1.4), Inches(0.22), Inches(0.35),
                  font_size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Calibri")

    # outcome banner
    rounded_rect(slide, Inches(0.25), Inches(5.68), Inches(12.8), Inches(0.65),
                 fill_color=RGBColor(0x00, 0x3A, 0x7A), corner_radius=0.08)
    txbox(slide, "✓  Final Dataset:  2,800 rows × 15 features  |  "
          "record_id excluded  |  Target: peak_load_risk (ordinal)",
          Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.48),
          font_size=11, bold=True, color=WHITE, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 5: SIMILARITY (Correlation) ─────────────────────────────────────────
def build_similarity_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "03 · Similarity", Inches(0.25), Inches(0.2))
    txbox(slide, "Feature Correlation Analysis",
          Inches(0.25), Inches(0.56), Inches(8.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    txbox(slide, "Pearson Correlation Heatmap — Numerical Features",
          Inches(0.25), Inches(1.35), Inches(7.0), Inches(0.35),
          font_size=10, color=GREY, italic=True, font_name="Calibri")

    add_image(slide, IMG["correlation"], Inches(0.25), Inches(1.75), Inches(7.6))

    # insight cards (right side)
    txbox(slide, "KEY INSIGHTS",
          Inches(8.25), Inches(1.35), Inches(4.8), Inches(0.35),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    insights = [
        (BLUE,      "Near-Zero Correlations",
                    "All Pearson coefficients fall close to 0 — features are largely independent."),
        (ORANGE,    "Non-Linear Risk Drivers",
                    "Weak linear correlations suggest peak load risk emerges from complex, non-linear interactions."),
        (TEAL,      "Ensemble Justified",
                    "Absence of dominant predictors supports the use of Random Forest over linear models."),
        (RED_WARN,  "Feature Set Limitation",
                    "Limited predictive power indicates external features (weather, pricing) are needed."),
    ]

    for i, (clr, title, body) in enumerate(insights):
        top = Inches(1.78) + i * Inches(1.35)
        rounded_rect(slide, Inches(8.25), top, Inches(4.85), Inches(1.2),
                     fill_color=CARD_BG, corner_radius=0.1)
        rect(slide, Inches(8.25), top, Inches(0.06), Inches(1.2), fill_color=clr)
        txbox(slide, title,
              Inches(8.42), top + Inches(0.1), Inches(4.5), Inches(0.38),
              font_size=11, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, body,
              Inches(8.42), top + Inches(0.48), Inches(4.5), Inches(0.65),
              font_size=9, color=GREY, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 6: DISSIMILARITY (Euclidean) ────────────────────────────────────────
def build_dissimilarity_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "03 · Dissimilarity", Inches(0.25), Inches(0.2))
    txbox(slide, "Euclidean Distance Across City Zones",
          Inches(0.25), Inches(0.56), Inches(8.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    add_image(slide, IMG["euclidean"], Inches(0.25), Inches(1.45), Inches(7.8))

    txbox(slide, "ZONE SIMILARITY MATRIX INSIGHTS",
          Inches(8.45), Inches(1.35), Inches(4.6), Inches(0.35),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    # most similar
    rounded_rect(slide, Inches(8.45), Inches(1.75), Inches(4.65), Inches(1.45),
                 fill_color=RGBColor(0x00, 0x25, 0x50), corner_radius=0.1)
    rect(slide, Inches(8.45), Inches(1.75), Inches(0.06), Inches(1.45), fill_color=GREEN_STAT)
    txbox(slide, "MOST SIMILAR",
          Inches(8.62), Inches(1.82), Inches(4.3), Inches(0.28),
          font_size=8, bold=True, color=GREEN_STAT, font_name="Calibri")
    txbox(slide, "South  ↔  West",
          Inches(8.62), Inches(2.08), Inches(4.3), Inches(0.5),
          font_size=20, bold=True, color=WHITE, font_name="Calibri")
    txbox(slide, "Distance ≈ 0.85  |  Comparable load profiles",
          Inches(8.62), Inches(2.56), Inches(4.3), Inches(0.5),
          font_size=9, color=GREY, font_name="Calibri")

    # most dissimilar
    rounded_rect(slide, Inches(8.45), Inches(3.35), Inches(4.65), Inches(1.45),
                 fill_color=RGBColor(0x3A, 0x10, 0x10), corner_radius=0.1)
    rect(slide, Inches(8.45), Inches(3.35), Inches(0.06), Inches(1.45), fill_color=RED_WARN)
    txbox(slide, "MOST DISSIMILAR",
          Inches(8.62), Inches(3.42), Inches(4.3), Inches(0.28),
          font_size=8, bold=True, color=RED_WARN, font_name="Calibri")
    txbox(slide, "East  ↔  South",
          Inches(8.62), Inches(3.68), Inches(4.3), Inches(0.5),
          font_size=20, bold=True, color=WHITE, font_name="Calibri")
    txbox(slide, "Distance ≈ 3.64  |  Fundamentally different behaviour",
          Inches(8.62), Inches(4.16), Inches(4.3), Inches(0.5),
          font_size=9, color=GREY, font_name="Calibri")

    # implication
    rounded_rect(slide, Inches(8.45), Inches(4.95), Inches(4.65), Inches(1.3),
                 fill_color=CARD_BG, corner_radius=0.1)
    txbox(slide, "💡  POLICY IMPLICATION",
          Inches(8.62), Inches(5.02), Inches(4.3), Inches(0.3),
          font_size=9, bold=True, color=HIGHLIGHT, font_name="Calibri")
    txbox(slide, "South-West strategies can be shared.\nEast zone requires a distinct, tailored grid plan.",
          Inches(8.62), Inches(5.32), Inches(4.3), Inches(0.78),
          font_size=9.5, color=OFFWHITE, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 7: CLASSIFICATION ────────────────────────────────────────────────────
def build_classification_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "04 · Classification", Inches(0.25), Inches(0.2))
    txbox(slide, "Random Forest Classifier — Results",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    # big accuracy hero
    rounded_rect(slide, Inches(0.25), Inches(1.45), Inches(3.5), Inches(2.2),
                 fill_color=RGBColor(0x00, 0x24, 0x60), corner_radius=0.12)
    rect(slide, Inches(0.25), Inches(1.45), Inches(3.5), Inches(0.07), fill_color=HIGHLIGHT)
    txbox(slide, "OVERALL\nACCURACY",
          Inches(0.4), Inches(1.58), Inches(3.2), Inches(0.55),
          font_size=10, bold=True, color=HIGHLIGHT, font_name="Calibri")
    txbox(slide, "45.18%",
          Inches(0.35), Inches(2.08), Inches(3.35), Inches(0.95),
          font_size=44, bold=True, color=WHITE, font_name="Calibri")
    txbox(slide, "with SMOTE rebalancing",
          Inches(0.4), Inches(2.98), Inches(3.2), Inches(0.4),
          font_size=9, color=GREY, italic=True, font_name="Calibri")

    # confusion matrix
    add_image(slide, IMG["confusion"], Inches(0.25), Inches(3.78), Inches(3.6))
    txbox(slide, "Confusion Matrix — Random Forest",
          Inches(0.25), Inches(7.0), Inches(3.6), Inches(0.3),
          font_size=7.5, italic=True, color=GREY, align=PP_ALIGN.CENTER, font_name="Calibri")

    # KPI cards (class-level metrics)
    metrics = [
        ("Low Risk",     "F1 = 0.58", "P: 0.55  R: 0.62  N=281", GREEN_STAT),
        ("Medium Risk",  "F1 = 0.33", "P: 0.34  R: 0.33  N=166", ORANGE),
        ("High Risk",    "F1 = 0.25", "P: 0.30  R: 0.21  N=113", RED_WARN),
        ("Weighted Avg", "F1 = 0.44", "P: 0.43  R: 0.45  N=560", GREY),
    ]

    kpi_top = Inches(1.45)
    kpi_left = Inches(4.1)
    for i, (cls, f1, detail, clr) in enumerate(metrics):
        kw = Inches(2.25)
        cx = kpi_left + i * (kw + Inches(0.18))
        rounded_rect(slide, cx, kpi_top, kw, Inches(1.65), fill_color=CARD_BG, corner_radius=0.1)
        rect(slide, cx, kpi_top, kw, Inches(0.06), fill_color=clr)
        txbox(slide, cls, cx + Inches(0.12), kpi_top + Inches(0.12), kw - Inches(0.24), Inches(0.35),
              font_size=9, bold=True, color=clr, font_name="Calibri")
        txbox(slide, f1, cx + Inches(0.1), kpi_top + Inches(0.48), kw - Inches(0.2), Inches(0.72),
              font_size=22, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, detail, cx + Inches(0.12), kpi_top + Inches(1.22), kw - Inches(0.24), Inches(0.35),
              font_size=8, color=GREY, font_name="Calibri")

    # methodology summary
    txbox(slide, "METHODOLOGY",
          Inches(4.1), Inches(3.3), Inches(9.0), Inches(0.32),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    method_points = [
        ("80/20 Split",        "Stratified train/test split preserving class proportions"),
        ("SMOTE",              "Synthetic Minority Over-sampling on training set only"),
        ("100 Trees",          "RandomForestClassifier — n_estimators=100, default params"),
        ("3-Class Target",     "Low (0) · Medium (1) · High (2) — ordinal risk levels"),
    ]
    for i, (tag, desc) in enumerate(method_points):
        row_left = Inches(4.1) + (i % 2) * Inches(4.55)
        row_top  = Inches(3.68) if i < 2 else Inches(4.72)
        rounded_rect(slide, row_left, row_top, Inches(4.3), Inches(0.85),
                     fill_color=CARD_BG, corner_radius=0.08)
        txbox(slide, tag, row_left + Inches(0.12), row_top + Inches(0.08),
              Inches(1.4), Inches(0.32), font_size=9.5, bold=True, color=BLUE, font_name="Calibri")
        txbox(slide, desc, row_left + Inches(0.12), row_top + Inches(0.4),
              Inches(4.0), Inches(0.36), font_size=8.5, color=GREY, font_name="Calibri")

    # interpretation
    txbox(slide, "INTERPRETATION",
          Inches(4.1), Inches(5.72), Inches(9.0), Inches(0.32),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    interp = ("Moderate accuracy reflects inherent complexity of 3-class risk prediction with weak linear feature "
              "separability. High-risk miss-rate (F1 = 0.25) is the critical operational gap — "
              "SMOTE improved minority class recall but non-linear boundaries limit achievable accuracy.")
    txbox(slide, interp, Inches(4.1), Inches(6.05), Inches(9.0), Inches(0.85),
          font_size=9.5, color=OFFWHITE, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 8: FEATURE IMPORTANCE ───────────────────────────────────────────────
def build_feature_importance_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "05 · Feature Importance", Inches(0.25), Inches(0.2))
    txbox(slide, "What Drives Peak Load Risk?",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    # large chart — hero element
    add_image(slide, IMG["feature_imp"], Inches(0.25), Inches(1.42), Inches(8.5))

    # callout cards right
    txbox(slide, "KEY TAKEAWAYS",
          Inches(9.1), Inches(1.42), Inches(4.0), Inches(0.32),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    callouts = [
        (BLUE,      "#1  Energy Dispensed (kWh)",
                    "Strongest predictor — total energy per session directly correlates with grid stress."),
        (TEAL,      "#2  Grid Load (MW)",
                    "Direct measure of instantaneous grid draw — naturally linked to risk level."),
        (ORANGE,    "#3  Renewable Energy %",
                    "Moderates grid stress — higher renewable share reduces peak risk contribution."),
        (HIGHLIGHT, "#4  No Single Dominant Feature",
                    "Uniform distribution of scores confirms multifactorial risk — ensemble model is the right choice."),
    ]
    for i, (clr, title, body) in enumerate(callouts):
        top = Inches(1.8) + i * Inches(1.38)
        rounded_rect(slide, Inches(9.1), top, Inches(4.0), Inches(1.22),
                     fill_color=CARD_BG, corner_radius=0.1)
        rect(slide, Inches(9.1), top, Inches(0.06), Inches(1.22), fill_color=clr)
        txbox(slide, title, Inches(9.27), top + Inches(0.1),
              Inches(3.75), Inches(0.38), font_size=10.5, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, body, Inches(9.27), top + Inches(0.5),
              Inches(3.75), Inches(0.65), font_size=8.5, color=GREY, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 9: CLASS DISTRIBUTION ───────────────────────────────────────────────
def build_class_dist_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "06 · Class Distribution", Inches(0.25), Inches(0.2))
    txbox(slide, "Target Variable: Peak Load Risk Distribution",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    # chart
    add_image(slide, IMG["class_dist"], Inches(0.25), Inches(1.45), Inches(8.1))

    # percentage cards
    txbox(slide, "CLASS BREAKDOWN",
          Inches(8.65), Inches(1.45), Inches(4.3), Inches(0.32),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    classes = [
        ("Low Risk",    "1,404", "50.1%", GREEN_STAT,
         "Majority class · Most reliably classified (F1=0.58)"),
        ("Medium Risk", "830",   "29.6%", ORANGE,
         "Transitional class · Hardest to classify (F1=0.33)"),
        ("High Risk",   "566",   "20.2%", RED_WARN,
         "Minority class · Operationally critical (F1=0.25)"),
    ]
    for i, (label, count, pct, clr, note) in enumerate(classes):
        top = Inches(1.88) + i * Inches(1.72)
        rounded_rect(slide, Inches(8.65), top, Inches(4.45), Inches(1.55),
                     fill_color=CARD_BG, corner_radius=0.1)
        rect(slide, Inches(8.65), top, Inches(0.06), Inches(1.55), fill_color=clr)
        txbox(slide, pct, Inches(8.82), top + Inches(0.1),
              Inches(1.4), Inches(0.72), font_size=32, bold=True, color=clr, font_name="Calibri")
        txbox(slide, f"{label}  ({count} records)",
              Inches(8.82), top + Inches(0.8), Inches(4.1), Inches(0.35),
              font_size=10, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, note,
              Inches(8.82), top + Inches(1.15), Inches(4.1), Inches(0.32),
              font_size=8, color=GREY, font_name="Calibri")

    # SMOTE imbalance note
    rounded_rect(slide, Inches(8.65), Inches(7.04), Inches(4.45), Inches(0.32),
                 fill_color=RGBColor(0x1A, 0x3A, 0x00), corner_radius=0.08)
    txbox(slide, "→  SMOTE applied on training set to address imbalance",
          Inches(8.75), Inches(7.07), Inches(4.2), Inches(0.26),
          font_size=8.5, color=GREEN_STAT, bold=True, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 10: K-MEANS CLUSTERING ──────────────────────────────────────────────
def build_kmeans_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "07 · Clustering", Inches(0.25), Inches(0.2))
    txbox(slide, "K-Means Clustering — 3 Operational Profiles",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    txbox(slide, "Algorithm: K-Means  |  k = 3  |  Features: energy_dispensed_kwh × grid_load_mw (scaled)",
          Inches(0.25), Inches(1.35), Inches(12.0), Inches(0.35),
          font_size=9.5, color=GREY, italic=True, font_name="Calibri")

    # cluster profile cards
    clusters = [
        ("Cluster 0",   "Residential\nOvernight",
         "Low energy · Low grid load\nSlow chargers · Extended duration",
         "⬇ Minimal grid impact",   RGBColor(0x26, 0x00, 0x8F)),   # purple

        ("Cluster 1",   "Urban\nFast Charger",
         "Moderate energy · Moderate load\nMix of commuter & commercial users",
         "↔ Standard grid operations", TEAL),

        ("Cluster 2",   "Commercial\nHub",
         "High energy · High grid load\nUltra-fast / multi-vehicle events",
         "⬆ Highest grid stress",    BLUE),
    ]
    card_w = Inches(4.0)
    for i, (cluster, name, desc, impact, clr) in enumerate(clusters):
        cx = Inches(0.25) + i * (card_w + Inches(0.18))
        cy = Inches(1.8)
        rounded_rect(slide, cx, cy, card_w, Inches(2.1), fill_color=CARD_BG, corner_radius=0.1)
        rect(slide, cx, cy, card_w, Inches(0.07), fill_color=clr)
        txbox(slide, cluster, cx + Inches(0.18), cy + Inches(0.12),
              card_w - Inches(0.36), Inches(0.3), font_size=9, bold=True, color=clr, font_name="Calibri")
        txbox(slide, name, cx + Inches(0.18), cy + Inches(0.42),
              card_w - Inches(0.36), Inches(0.65), font_size=16, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, desc, cx + Inches(0.18), cy + Inches(1.08),
              card_w - Inches(0.36), Inches(0.62), font_size=9, color=GREY, font_name="Calibri")
        rounded_rect(slide, cx + Inches(0.18), cy + Inches(1.72),
                     card_w - Inches(0.36), Inches(0.28), fill_color=clr, corner_radius=0.15)
        txbox(slide, impact, cx + Inches(0.22), cy + Inches(1.74),
              card_w - Inches(0.44), Inches(0.25), font_size=8, bold=True, color=WHITE, font_name="Calibri")

    # cluster scatter (from dashboard panel 2)
    txbox(slide, "Cluster Scatter Plot (from Final Dashboard — Panel 2)",
          Inches(0.25), Inches(4.1), Inches(10.0), Inches(0.32),
          font_size=9, color=GREY, italic=True, font_name="Calibri")
    add_image(slide, IMG["dashboard"], Inches(0.25), Inches(4.42), Inches(12.8), Inches(2.7))

    # operational insight
    rounded_rect(slide, Inches(0.25), Inches(7.08), Inches(12.8), Inches(0.3),
                 fill_color=RGBColor(0x00, 0x3A, 0x5A), corner_radius=0.06)
    txbox(slide, "💡  Cluster-specific grid policies outperform uniform limits — "
          "tailor demand-response triggers to each operational archetype",
          Inches(0.45), Inches(7.1), Inches(12.2), Inches(0.26),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 11: OUTLIER DETECTION ───────────────────────────────────────────────
def build_outlier_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "08 · Anomaly Detection", Inches(0.25), Inches(0.2))
    txbox(slide, "Isolation Forest — Anomaly Detection",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    # big anomaly count hero
    rounded_rect(slide, Inches(0.25), Inches(1.45), Inches(3.4), Inches(2.5),
                 fill_color=RGBColor(0x35, 0x08, 0x18), corner_radius=0.12)
    rect(slide, Inches(0.25), Inches(1.45), Inches(3.4), Inches(0.07), fill_color=RED_WARN)
    txbox(slide, "⚠  ANOMALIES\nFLAGGED",
          Inches(0.45), Inches(1.6), Inches(3.0), Inches(0.65),
          font_size=10, bold=True, color=RED_WARN, font_name="Calibri")
    txbox(slide, "84",
          Inches(0.4), Inches(2.2), Inches(3.2), Inches(0.95),
          font_size=72, bold=True, color=WHITE, font_name="Calibri")
    txbox(slide, "out of 2,800 sessions\n3% contamination rate",
          Inches(0.45), Inches(3.12), Inches(3.0), Inches(0.65),
          font_size=9.5, color=GREY, font_name="Calibri")

    # normal vs anomaly
    rounded_rect(slide, Inches(0.25), Inches(4.15), Inches(1.6), Inches(1.0),
                 fill_color=CARD_BG, corner_radius=0.1)
    txbox(slide, "2,716", Inches(0.38), Inches(4.22), Inches(1.35), Inches(0.5),
          font_size=22, bold=True, color=GREY, font_name="Calibri")
    txbox(slide, "Normal\nSessions", Inches(0.38), Inches(4.68), Inches(1.35), Inches(0.38),
          font_size=8, color=GREY, font_name="Calibri")

    rounded_rect(slide, Inches(2.0), Inches(4.15), Inches(1.6), Inches(1.0),
                 fill_color=RGBColor(0x35, 0x08, 0x18), corner_radius=0.1)
    txbox(slide, "84", Inches(2.13), Inches(4.22), Inches(1.35), Inches(0.5),
          font_size=22, bold=True, color=RED_WARN, font_name="Calibri")
    txbox(slide, "Anomalous\nSessions", Inches(2.13), Inches(4.68), Inches(1.35), Inches(0.38),
          font_size=8, color=RED_WARN, font_name="Calibri")

    # chart (panel 3 of dashboard)
    txbox(slide, "Anomaly Overlay (from Dashboard — Panel 3: Red = Anomalies)",
          Inches(3.9), Inches(1.35), Inches(9.2), Inches(0.32),
          font_size=9, color=GREY, italic=True, font_name="Calibri")
    add_image(slide, IMG["dashboard"], Inches(3.85), Inches(1.68), Inches(9.25), Inches(3.85))

    # causes
    txbox(slide, "POTENTIAL CAUSES",
          Inches(0.25), Inches(5.35), Inches(3.5), Inches(0.3),
          font_size=9, bold=True, color=TEAL, font_name="Calibri")

    causes = [
        (RED_WARN, "Hardware malfunction — erroneous power draw at specific stations"),
        (ORANGE,   "Grid baseline stress — simultaneous large-scale charging events"),
        (HIGHLIGHT, "Sensor faults — data recording anomalies at station level"),
    ]
    for i, (clr, txt) in enumerate(causes):
        top = Inches(5.7) + i * Inches(0.55)
        rect(slide, Inches(0.25), top + Inches(0.1), Inches(0.06), Inches(0.3), fill_color=clr)
        txbox(slide, txt, Inches(0.42), top + Inches(0.06), Inches(3.3), Inches(0.38),
              font_size=9, color=OFFWHITE, font_name="Calibri")

    # action panel
    rounded_rect(slide, Inches(3.85), Inches(5.65), Inches(9.25), Inches(1.68),
                 fill_color=CARD_BG, corner_radius=0.1)
    rect(slide, Inches(3.85), Inches(5.65), Inches(9.25), Inches(0.07), fill_color=RED_WARN)
    txbox(slide, "OPERATIONAL RESPONSE ACTIONS",
          Inches(4.05), Inches(5.76), Inches(9.0), Inches(0.32),
          font_size=9, bold=True, color=RED_WARN, font_name="Calibri")

    actions = [
        "Investigate flagged station records for hardware faults",
        "Pre-position emergency battery buffer reserves",
        "Implement real-time circuit breaker protocols for similar future events",
    ]
    for i, act in enumerate(actions):
        txbox(slide, f"→  {act}",
              Inches(4.05), Inches(6.12) + i * Inches(0.38), Inches(8.8), Inches(0.34),
              font_size=9.5, color=OFFWHITE, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 12: DASHBOARD ────────────────────────────────────────────────────────
def build_dashboard_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "09 · Dashboard", Inches(0.25), Inches(0.2))
    txbox(slide, "Three-Panel Executive Dashboard",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")

    panels = [
        (ORANGE,    "Panel 1",  "Temporal Risk",   "Hourly grid load volatility"),
        (TEAL,      "Panel 2",  "K-Means Clusters","Operational zone segmentation"),
        (RED_WARN,  "Panel 3",  "Anomaly Detection","84 flagged events in red"),
    ]
    pw = Inches(4.22)
    for i, (clr, tag, title, sub) in enumerate(panels):
        cx = Inches(0.25) + i * (pw + Inches(0.1))
        txbox(slide, f"{tag}: {title}",
              cx, Inches(1.22), pw, Inches(0.28),
              font_size=9, bold=True, color=clr, font_name="Calibri")

    # dashboard image — hero
    add_image(slide, IMG["dashboard"], Inches(0.15), Inches(1.55), Inches(13.0), Inches(4.65))

    # panel subtitles below
    for i, (clr, tag, title, sub) in enumerate(panels):
        cx = Inches(0.25) + i * (pw + Inches(0.1))
        rounded_rect(slide, cx, Inches(6.28), pw, Inches(0.62),
                     fill_color=CARD_BG, corner_radius=0.08)
        rect(slide, cx, Inches(6.28), pw, Inches(0.06), fill_color=clr)
        txbox(slide, sub, cx + Inches(0.12), Inches(6.38), pw - Inches(0.24), Inches(0.42),
              font_size=9.5, color=OFFWHITE, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 13: CONCLUSION ───────────────────────────────────────────────────────
def build_conclusion_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "10 · Conclusions", Inches(0.25), Inches(0.2))
    txbox(slide, "Key Findings & Takeaways",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    findings = [
        (TEAL,      "Clean Dataset",
                    "2,800 records · 10 attributes · 0 missing values · 0 duplicates — no imputation needed"),
        (BLUE,      "Weak Feature Correlations",
                    "Pearson near-zero → non-linear risk drivers · Zone dissimilarity supports zone-specific grid policy"),
        (ORANGE,    "RF Accuracy: 45.18%",
                    "Moderate but informative · SMOTE improved minority recall · High-risk F1=0.25 is the critical gap"),
        (HIGHLIGHT, "Multifactorial Risk",
                    "energy_kwh, grid_load_mw, renewable_% are top 3 predictors · No single dominant feature"),
        (GREEN_STAT,"3 Operational Clusters",
                    "K-Means: Residential · Urban Fast · Commercial Hub → cluster-specific demand-response policies"),
        (RED_WARN,  "84 Anomalous Events",
                    "Isolation Forest (3% contamination) · Black-swan grid stress moments requiring investigation"),
    ]

    cw = Inches(4.2)
    for i, (clr, title, body) in enumerate(findings):
        col = i % 3
        row = i // 3
        cx = Inches(0.25) + col * (cw + Inches(0.1))
        cy = Inches(1.48) + row * Inches(2.28)
        rounded_rect(slide, cx, cy, cw, Inches(2.1), fill_color=CARD_BG, corner_radius=0.12)
        rect(slide, cx, cy, cw, Inches(0.07), fill_color=clr)
        # icon-like number
        rounded_rect(slide, cx + Inches(0.15), cy + Inches(0.18),
                     Inches(0.42), Inches(0.42), fill_color=clr, corner_radius=0.5)
        txbox(slide, str(i + 1),
              cx + Inches(0.15), cy + Inches(0.2), Inches(0.42), Inches(0.38),
              font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
        txbox(slide, title,
              cx + Inches(0.65), cy + Inches(0.18), cw - Inches(0.82), Inches(0.42),
              font_size=12, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, body,
              cx + Inches(0.15), cy + Inches(0.72), cw - Inches(0.3), Inches(1.28),
              font_size=9, color=GREY, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 14: FUTURE WORK ──────────────────────────────────────────────────────
def build_future_work_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)
    accent_bar(slide)

    section_pill(slide, "11 · Future Work", Inches(0.25), Inches(0.2))
    txbox(slide, "Recommendations & Next Steps",
          Inches(0.25), Inches(0.56), Inches(10.0), Inches(0.65),
          font_size=30, bold=True, color=WHITE, font_name="Calibri")
    divider_line(slide, Inches(0.25), Inches(1.26), Inches(12.8))

    future = [
        (BLUE,      "External Feature Enrichment",
                    "Incorporate ambient temperature, baseline neighbourhood power consumption, "
                    "and real-time electricity pricing to improve RF accuracy beyond 45%."),
        (TEAL,      "Advanced Model Architecture",
                    "Explore Gradient Boosting (XGBoost/LightGBM), Neural Networks, and "
                    "class-weight tuning as alternatives to SMOTE for imbalance."),
        (ORANGE,    "Dynamic Time-of-Use Pricing",
                    "Grid operators should implement cheaper off-peak rates to flatten the "
                    "hourly demand curve and reduce infrastructure degradation."),
        (GREEN_STAT,"Zone-Specific Forecasting",
                    "Build separate predictive models per city zone — especially East zone "
                    "which is operationally distinct from all others (max distance ≈ 3.64)."),
        (HIGHLIGHT, "Real-Time Anomaly Monitoring",
                    "Deploy Isolation Forest as a live streaming pipeline to flag "
                    "anomalous events in real-time and trigger automated circuit responses."),
        (RED_WARN,  "Dataset Expansion",
                    "Merge with city-wide macroeconomic, weather, and grid telemetry datasets "
                    "for a more generalisable model beyond synthetic Kaggle data."),
    ]

    cw = Inches(4.2)
    for i, (clr, title, body) in enumerate(future):
        col = i % 3
        row = i // 3
        cx = Inches(0.25) + col * (cw + Inches(0.1))
        cy = Inches(1.48) + row * Inches(2.6)
        rounded_rect(slide, cx, cy, cw, Inches(2.45), fill_color=CARD_BG, corner_radius=0.12)
        rect(slide, cx, cy, cw, Inches(0.07), fill_color=clr)
        rounded_rect(slide, cx + Inches(0.15), cy + Inches(0.18),
                     Inches(0.42), Inches(0.42), fill_color=clr, corner_radius=0.5)
        txbox(slide, "→",
              cx + Inches(0.15), cy + Inches(0.2), Inches(0.42), Inches(0.38),
              font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
        txbox(slide, title,
              cx + Inches(0.65), cy + Inches(0.18), cw - Inches(0.82), Inches(0.42),
              font_size=11.5, bold=True, color=WHITE, font_name="Calibri")
        txbox(slide, body,
              cx + Inches(0.15), cy + Inches(0.72), cw - Inches(0.3), Inches(1.62),
              font_size=9, color=GREY, font_name="Calibri")

    slide_number_tag(slide, n, total)
    return slide


# ── SLIDE 15: THANK YOU ────────────────────────────────────────────────────────
def build_thankyou_slide(prs, n, total):
    slide = add_blank_slide(prs)
    fill_slide_bg(slide, NAVY)

    rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill_color=BLUE)
    rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), fill_color=TEAL)
    rect(slide, 0, 0, Inches(0.06), SLIDE_H, fill_color=TEAL)
    rect(slide, SLIDE_W - Inches(0.06), 0, Inches(0.06), SLIDE_H, fill_color=BLUE)

    txbox(slide, "Thank You",
          Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.8),
          font_size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

    divider_line(slide, Inches(3.0), Inches(3.35), Inches(7.33), color=BLUE, thickness_pt=2)

    txbox(slide, "CE634 · Data Mining and Data Warehousing",
          Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.5),
          font_size=14, color=TEAL, align=PP_ALIGN.CENTER, font_name="Calibri")

    txbox(slide, "Dewpearl Gonsalves  ·  Shail Joshi  ·  Flyson Dias",
          Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.45),
          font_size=12, color=GREY, align=PP_ALIGN.CENTER, font_name="Calibri")

    txbox(slide, "Padre Conceicao College of Engineering  |  Academic Year 2025–2026",
          Inches(1.5), Inches(4.75), Inches(10.3), Inches(0.4),
          font_size=10, color=GREY, align=PP_ALIGN.CENTER, font_name="Calibri")

    # mini stats reminder
    stats_row = [
        ("2,800 Records", BLUE),
        ("45.18% RF Accuracy", ORANGE),
        ("84 Anomalies", RED_WARN),
        ("3 Clusters", TEAL),
    ]
    row_w = Inches(2.9)
    for i, (label, clr) in enumerate(stats_row):
        cx = Inches(0.55) + i * (row_w + Inches(0.15))
        rounded_rect(slide, cx, Inches(5.75), row_w, Inches(0.65), fill_color=CARD_BG, corner_radius=0.12)
        txbox(slide, label, cx + Inches(0.1), Inches(5.87), row_w - Inches(0.2), Inches(0.42),
              font_size=11, bold=True, color=clr, align=PP_ALIGN.CENTER, font_name="Calibri")

    txbox(slide, "Open for Questions",
          Inches(1.5), Inches(6.7), Inches(10.3), Inches(0.55),
          font_size=14, color=GREY, align=PP_ALIGN.CENTER, italic=True, font_name="Calibri")

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build_presentation():
    prs = new_prs()
    TOTAL = 15

    print("Building CE634 Premium Viva Presentation...")

    print("  Slide 01 — Title")
    build_title_slide(prs, 1, TOTAL)

    print("  Slide 02 — Agenda")
    build_agenda_slide(prs, 2, TOTAL)

    print("  Slide 03 — Dataset")
    build_dataset_slide(prs, 3, TOTAL)

    print("  Slide 04 — Preprocessing Pipeline")
    build_preprocessing_slide(prs, 4, TOTAL)

    print("  Slide 05 — Similarity (Correlation)")
    build_similarity_slide(prs, 5, TOTAL)

    print("  Slide 06 — Dissimilarity (Euclidean)")
    build_dissimilarity_slide(prs, 6, TOTAL)

    print("  Slide 07 — Classification (Random Forest)")
    build_classification_slide(prs, 7, TOTAL)

    print("  Slide 08 — Feature Importance")
    build_feature_importance_slide(prs, 8, TOTAL)

    print("  Slide 09 — Class Distribution")
    build_class_dist_slide(prs, 9, TOTAL)

    print("  Slide 10 — K-Means Clustering")
    build_kmeans_slide(prs, 10, TOTAL)

    print("  Slide 11 — Outlier Detection")
    build_outlier_slide(prs, 11, TOTAL)

    print("  Slide 12 — Dashboard")
    build_dashboard_slide(prs, 12, TOTAL)

    print("  Slide 13 — Conclusion")
    build_conclusion_slide(prs, 13, TOTAL)

    print("  Slide 14 — Future Work")
    build_future_work_slide(prs, 14, TOTAL)

    print("  Slide 15 — Thank You")
    build_thankyou_slide(prs, 15, TOTAL)

    prs.save(OUTPUT)
    print(f"\n[OK] Presentation saved:\n   {OUTPUT}\n")
    print(f"Total slides: {TOTAL}")


if __name__ == "__main__":
    build_presentation()
